# -*- coding: utf-8 -*-
"""Generate PowerPoint from HTML - Salamander Regeneration Presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# Unicode helpers for Chinese curly quotes
LQ = '“'  # "
RQ = '”'  # "

# --- Color Constants ---
BG_DEEP = RGBColor(0x0A, 0x16, 0x28)
BG_SURFACE = RGBColor(0x12, 0x23, 0x40)
BG_SURFACE2 = RGBColor(0x16, 0x2D, 0x52)
OUC_BLUE = RGBColor(0x00, 0x3B, 0x7B)
OUC_CYAN = RGBColor(0x00, 0x88, 0xCC)
CORAL = RGBColor(0xE8, 0x61, 0x4D)
GOLD = RGBColor(0xC9, 0xA9, 0x6E)
TEXT_WHITE = RGBColor(0xE8, 0xED, 0xF2)
TEXT_MUTED = RGBColor(0x8A, 0x9B, 0xB5)
TEXT_MUTED2 = RGBColor(0x5A, 0x6D, 0x88)
BORDER = RGBColor(0x8A, 0x9B, 0xB5)
SEMI_BLUE = RGBColor(0x00, 0x3B, 0x7B)

IMG_DIR = "c:/Users/流光/Desktop/pptx/extracted_images"
OUTPUT = "c:/Users/流光/Desktop/pptx/蝾螈再生术_OUC版.pptx"

FONT_SERIF = 'Noto Serif SC'
FONT_SANS = 'Noto Sans SC'
FONT_MONO = 'Courier New'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def S():
    """Create a new slide with dark bg."""
    s = prs.slides.add_slide(BLANK)
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DEEP
    return s


def TB(s, l, t, w, h, txt='', fs=14, fc=TEXT_WHITE, b=False, fn=FONT_SANS,
       al=PP_ALIGN.LEFT, it=False):
    """Simple text box."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(fs)
    p.font.color.rgb = fc
    p.font.bold = b
    p.font.name = fn
    p.font.italic = it
    p.alignment = al
    return tf


def RT(s, l, t, w, h):
    """Rich text frame."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def AP(tf, txt='', fs=14, fc=TEXT_WHITE, b=False, fn=FONT_SANS, it=False,
       al=PP_ALIGN.LEFT, sb=0, sa=0):
    """Add paragraph to text frame."""
    p = tf.add_paragraph()
    p.text = txt
    p.font.size = Pt(fs)
    p.font.color.rgb = fc
    p.font.bold = b
    p.font.name = fn
    p.font.italic = it
    p.alignment = al
    if sb:
        p.space_before = Pt(sb)
    if sa:
        p.space_after = Pt(sa)
    return p


def RE(s, l, t, w, h, fill=None, border=None, bw=None):
    """Rectangle shape."""
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                            Inches(w), Inches(h))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border
        if bw:
            sh.line.width = Pt(bw)
    else:
        sh.line.fill.background()
    return sh


def AC(s, l, t):
    """Accent rule line."""
    RE(s, l, t, 0.67, 0.03, fill=OUC_CYAN)


def WD(s):
    """Wave divider at bottom."""
    RE(s, 0, 7.45, 13.333, 0.04, fill=OUC_CYAN)


def WM(s):
    """OUC watermark top-right."""
    TB(s, 9.5, 0.35, 3.6, 0.35,
       '中国海洋大学  OCEAN UNIVERSITY OF CHINA',
       fs=8, fc=TEXT_MUTED, fn=FONT_SANS, al=PP_ALIGN.RIGHT)


def IMG(s, name, l, t, w, h=None):
    """Add image if exists."""
    p = os.path.join(IMG_DIR, name)
    if os.path.exists(p):
        if h:
            return s.shapes.add_picture(p, Inches(l), Inches(t), Inches(w), Inches(h))
        return s.shapes.add_picture(p, Inches(l), Inches(t), Inches(w))
    return None


def TAG(s, l, t, txt):
    """Simple tag label."""
    TB(s, l, t, 2.5, 0.3, txt, fs=9, fc=OUC_CYAN, fn=FONT_MONO)


# =====================================================================
# SLIDE 1: COVER
# =====================================================================
s = S()
WM(s); WD(s)
# Left column
RE(s, 1.11, 1.5, 0.04, 0.5, fill=OUC_BLUE)
TB(s, 1.35, 1.45, 4.5, 0.28, '中国海洋大学 · 海洋生命学院',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
TB(s, 1.35, 1.72, 4.5, 0.2,
   'OCEAN UNIVERSITY OF CHINA · COLLEGE OF MARINE LIFE SCIENCES',
   fs=7, fc=TEXT_MUTED2, fn=FONT_MONO)
AC(s, 1.35, 2.05)
TB(s, 1.35, 2.2, 6.5, 2.2, '蝾螈的再生术',
   fs=72, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
TB(s, 1.35, 4.2, 5.5, 0.8,
   '解码自然界最精密的再生程序\n从细胞重编程到再生医学转化的前沿探索',
   fs=14, fc=TEXT_MUTED, fn=FONT_SANS)
# Tags
for i, tag in enumerate(['再生生物学', '基因调控网络',
                          '再生医学', '进化发育生物学']):
    RE(s, 1.35 + i * 2.0, 5.1, 1.7, 0.32, fill=SEMI_BLUE, border=OUC_CYAN, bw=0.5)
    TB(s, 1.45 + i * 2.0, 5.12, 1.5, 0.28, tag, fs=9, fc=OUC_CYAN, fn=FONT_MONO, al=PP_ALIGN.CENTER)
TB(s, 1.35, 5.55, 5.5, 0.25,
   '引用文献: Nature · Science · Cell · 中国海洋大学学报 (2020–2026)',
   fs=8, fc=TEXT_MUTED, fn=FONT_MONO)
# Right: image + stats
IMG(s, 'slide2_img3.jpeg', 7.8, 0.7, 4.5, 4.2)
stats = [('32Gb', '基因组大小\n人类10倍', OUC_CYAN),
         ('200+', '再生周期(天)\n完美重生', CORAL),
         ('5次+', '重复再生\n不降质', GOLD)]
for i, (v, lb, c) in enumerate(stats):
    x = 7.8 + i * 1.65
    TB(s, x, 5.2, 1.5, 0.8, v, fs=32, fc=c, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
    TB(s, x, 5.65, 1.5, 0.6, lb, fs=9, fc=TEXT_MUTED2, fn=FONT_MONO, al=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 2: Regeneration Overview
# =====================================================================
s = S(); WM(s); WD(s)
TB(s, 2.5, 0.8, 8.3, 0.28, '第一章 · 再生现象',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO, al=PP_ALIGN.CENTER)
AC(s, 6.3, 1.25)
TB(s, 1.5, 1.4, 10.3, 1.5,
   '再生能力\n进化树上的非凡分布',
   fs=48, fc=TEXT_WHITE, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
TB(s, 2.5, 2.85, 8.3, 0.6,
   f'再生不是{LQ}全或无{RQ}的性状 —— 在动物界中，它呈现从组织修复到完整个体重建的连续谱。\n理解这一谱系的分子基础，是再生生物学的核心命题。',
   fs=13, fc=TEXT_MUTED, fn=FONT_SANS, al=PP_ALIGN.CENTER)
# Species table
RE(s, 0.8, 3.6, 11.7, 1.8, fill=BG_SURFACE, border=BORDER, bw=0.5)
species = [
    ('\U0001f98e', '蝾螈', '肢体·心脏·视网膜\n脊髓·大脑·晶状体'),
    ('\U0001f41f', '斑马鱼', '心脏·鳍·视网膜\n脊髓·视神经'),
    ('\U0001fab1', '涡虫', '全身任何部位\n含中枢神经系统'),
    ('\U0001f42d', '小鼠', '趾尖·肝脏\n新生早期心脏'),
    ('\U0001f9ec', '人类', '肝脏·皮肤·指尖\n补偿性增生为主'),
]
for i, (icon, name, desc) in enumerate(species):
    x = 1.0 + i * 2.35
    TB(s, x, 3.8, 2.1, 0.5, icon, fs=24, fc=TEXT_WHITE, fn=FONT_SANS, al=PP_ALIGN.CENTER)
    TB(s, x, 4.2, 2.1, 0.38, name, fs=18, fc=OUC_CYAN, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
    TB(s, x, 4.55, 2.1, 0.7, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS, al=PP_ALIGN.CENTER)
    if i < 4:
        RE(s, x + 2.2, 3.8, 0.012, 1.4, fill=BORDER)
# Quote
RE(s, 2.0, 5.65, 9.3, 0.95, fill=RGBColor(0x00, 0x2D, 0x5E), border=OUC_CYAN, bw=1.5)
TB(s, 2.2, 5.72, 8.9, 0.5,
   f'{LQ}再生能力的进化丢失并非不可逆。我们正在学习如何唤醒哺乳动物中沉睡的再生程序。{RQ}',
   fs=14, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
TB(s, 2.2, 6.28, 8.9, 0.25,
   '—— Tanaka E M, Cell, 2016; 中国海洋大学发育与再生实验室, 2024',
   fs=9, fc=TEXT_MUTED2, fn=FONT_MONO)


# =====================================================================
# SLIDE 3: Salamander - King of Regeneration
# =====================================================================
s = S(); WM(s); WD(s)
# Left
TB(s, 1.0, 0.8, 4.5, 0.28, '第二章 · 模式物种',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
AC(s, 1.0, 1.15)
TB(s, 1.0, 1.3, 4.5, 1.3,
   '蝾螈\n脊椎动物再生之巅',
   fs=40, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
TB(s, 1.0, 2.8, 4.5, 1.0,
   '美西钝口螈（Ambystoma mexicanum）隶属有尾目钝口螈科，原产于墨西哥霍奇米尔科湖。它是目前唯一已知的、能够终生保持完美再生能力的四足脊椎动物模型。',
   fs=12, fc=TEXT_MUTED, fn=FONT_SANS)
for i, tag in enumerate(['基因组 32 Gb', '5次以上重复再生',
                          '精确位置记忆', '无瘢痕愈合']):
    TB(s, 1.0 + (i % 2) * 2.5, 4.0 + (i // 2) * 0.45, 2.3, 0.35, tag,
       fs=9, fc=OUC_CYAN, fn=FONT_MONO, al=PP_ALIGN.CENTER)
TB(s, 1.0, 5.2, 4.5, 0.25,
   'Nowoshilow S et al., Nature, 2018 · 中国海洋大学水生生物模式平台',
   fs=8, fc=TEXT_MUTED, fn=FONT_MONO)
# Right
IMG(s, 'slide1_img1.jpeg', 6.0, 0.7, 6.5, 3.5)
cards3 = [
    ('肢体再生', '约200天完成完整功能性肢体再生，骨骼·肌肉·神经·血管精确重建。', OUC_CYAN),
    ('脊髓再生', '完全横切后通过室管膜胶质细胞增殖和轴突再生实现功能恢复。', CORAL),
    ('心脏再生', '心室切除后无瘤痕修复，心肌细胞去分化增殖填补缺损。', GOLD),
    ('视网膜再生', '色素上皮细胞转分化产生完整神经视网膜，功能完全恢复。', OUC_BLUE),
]
for i, (title, desc, accent) in enumerate(cards3):
    x = 6.0 + (i % 2) * 3.45
    y = 4.4 + (i // 2) * 1.35
    RE(s, x, y, 3.2, 1.2, fill=BG_SURFACE, border=BORDER, bw=0.5)
    TB(s, x + 0.2, y + 0.1, 2.8, 0.32, title, fs=16, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, x + 0.2, y + 0.5, 2.8, 0.6, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS)


# =====================================================================
# SLIDE 4: Five Stages of Limb Regeneration
# =====================================================================
s = S(); WM(s); WD(s)
TB(s, 0.9, 0.5, 5.0, 0.28, '第三章 · 细胞机制',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
TB(s, 0.9, 0.9, 8.0, 0.65, '肢体再生的五个精密阶段',
   fs=36, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
TB(s, 0.9, 1.5, 11.5, 0.45,
   f'蝾螈肢体再生不是简单的{LQ}重新长出{RQ}，而是一个高度有序的多阶段细胞重编程过程。去分化而非成体干细胞的动员，是其区别于哺乳动物修复的核心特征。',
   fs=11, fc=TEXT_MUTED, fn=FONT_SANS)
stages = [
    ('01', '伤口愈合', '0–12 h', '表皮细胞快速迁移\n形成顶端外胚层帽(AEC)\n无纤维化瘢痕', OUC_CYAN),
    ('02', '去分化', '1–3 d', '成熟细胞退出分化\n重入细胞周期\n获得祖细胞多能性', CORAL),
    ('03', '芽基形成', '3–7 d', '去分化细胞聚集\n形成多能祖细胞团\nAEC-FGF 维持增殖', GOLD),
    ('04', '再分化', '21–50 d', '祖细胞重新分化\n骨骼·软骨·肌肉\n神经·血管同时构建', OUC_BLUE),
    ('05', '形态建成', '50–200 d', '三维结构精确重构\n模式形成+生长\n功能性肢体完美重生', CORAL),
]
for i, (num, title, time, desc, accent) in enumerate(stages):
    x = 0.4 + i * 2.55
    RE(s, x, 2.3, 2.3, 3.2, fill=BG_SURFACE, border=accent, bw=2)
    TB(s, x, 2.5, 2.3, 0.75, num, fs=40, fc=accent, fn=FONT_MONO, al=PP_ALIGN.CENTER)
    TB(s, x, 3.15, 2.3, 0.38, title, fs=18, fc=TEXT_WHITE, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
    TB(s, x, 3.5, 2.3, 0.28, time, fs=9, fc=TEXT_MUTED2, fn=FONT_MONO, al=PP_ALIGN.CENTER)
    TB(s, x + 0.15, 3.8, 2.0, 1.4, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS, al=PP_ALIGN.CENTER)
    if i < 4:
        TB(s, x + 2.3, 3.5, 0.3, 0.45, '▸', fs=20, fc=OUC_CYAN, fn=FONT_SANS, al=PP_ALIGN.CENTER)
IMG(s, 'slide4_img3.png', 0.9, 5.75, 11.5, 1.3)


# =====================================================================
# SLIDE 5: Blastema - Signal Hub
# =====================================================================
s = S(); WM(s); WD(s)
RE(s, 0, 0, 5.8, 7.5, fill=BG_SURFACE)
TB(s, 0.8, 0.6, 4.5, 0.28, '第三章 · 深度',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
AC(s, 0.8, 0.95)
TB(s, 0.8, 1.1, 4.5, 1.0, '芽基\n再生的信号枢纽',
   fs=34, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
b_items = [
    ('01', '细胞来源异质性',
     f'肌肉卫星细胞、成纤维细胞、施万细胞、软骨细胞协同去分化，形成混合祖细胞群体。谱系追踪证实：芽基细胞保留来源记忆但具跨谱系分化潜能。'),
    ('02', '位置信息编码',
     f'芽基细胞{LQ}记住{RQ}沿肢体近-远端轴的位置身份。HoxA 基因簇的共线性表达和视黄酸（RA）浓度梯度共同构成这一位置记忆的分子基础。'),
    ('03', '免疫豁免微环境',
     f'M2 型巨噬细胞和 Treg 细胞富集于芽基，分泌 IL-10、TGF-β3 等抗炎因子，构建{LQ}再生许可{RQ}微环境 —— 蝾螈与哺乳动物瘤痕愈合的关键分岔点。'),
]
for i, (num, title, desc) in enumerate(b_items):
    y = 2.3 + i * 1.55
    TB(s, 0.8, y, 0.45, 0.45, num, fs=22, fc=OUC_CYAN, fn=FONT_MONO)
    TB(s, 1.4, y, 4.0, 0.32, title, fs=18, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 1.4, y + 0.35, 4.0, 1.0, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS)
TB(s, 0.8, 7.0, 4.5, 0.25,
   'Kragl M et al., Nature, 2009 · Gerber T et al., Science, 2018',
   fs=8, fc=TEXT_MUTED2, fn=FONT_MONO)
# Right
RE(s, 5.8, 0, 7.533, 7.5, fill=BG_SURFACE2)
TB(s, 6.3, 0.6, 6.5, 0.28, '五大核心信号通路协同',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
pathways = [
    ('FGF 信号', 'AEC 分泌 FGF8/10，通过 FGFR1/2 维持芽基祖细胞增殖与未分化状态。阻断 FGF 信号导致再生中止。', OUC_CYAN),
    ('BMP 信号', 'BMP2/7 调控凋亡和模式形成。BMP 拮抗剂 Noggin 在 AEC 中高表达，精细控制 BMP 活性梯度。', CORAL),
    ('Wnt/β-catenin', '再生启动的核心开关。Wnt7a 决定前后轴极性，β-catenin 核转位是芽基形成的充分必要条件。', GOLD),
    ('Shh 信号', '芽基后端 ZPA 区分泌 Shh，建立前后轴模式，决定指趾身份。Shh 表达位置精确对应截肢水平。', OUC_BLUE),
    ('Notch 信号', 'Notch1-Delta 维持祖细胞状态，抑制过早分化。Notch 与 FGF 形成正反馈环，稳定芽基细胞命运。', OUC_CYAN),
]
for i, (title, desc, accent) in enumerate(pathways):
    y = 1.2 + i * 1.15
    RE(s, 6.3, y, 6.5, 1.0, fill=BG_SURFACE, border=BORDER, bw=0.5)
    TB(s, 6.5, y + 0.08, 6.1, 0.32, title, fs=16, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 6.5, y + 0.42, 6.1, 0.5, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS)


# =====================================================================
# SLIDE 6: Gene Regulatory Networks
# =====================================================================
s = S(); WM(s); WD(s)
TB(s, 0.8, 0.4, 5.0, 0.28, '第四章 · 分子密码',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
TB(s, 0.8, 0.72, 10.0, 0.65, '基因调控网络 再生的分子基础',
   fs=36, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
TB(s, 0.8, 1.32, 11.7, 0.48,
   '蝾螈再生涉及 200+ 个基因的协同时序表达。约30%的再生相关基因为蝾螈特有或功能分化。表观遗传重编程（尤其是 H3K27me3 去甲基化）是再生启动的先决条件。',
   fs=11, fc=TEXT_MUTED, fn=FONT_SANS)
# Left: Activators
TB(s, 0.8, 1.85, 5.5, 0.28, '再生激活因子',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
activators = [
    ('Msx1', '芽基形成关键转录因子，抑制肌源性分化程序（MyoD 下调），维持祖细胞多能性。', OUC_CYAN),
    ('Pax7', '肌肉卫星细胞激活与去分化的核心因子。Pax7+ 细胞是芽基中肌源性祖细胞的主要来源。', CORAL),
    ('Fgf8 / Fgf10', 'AEC 分泌的生长因子，维持芽基祖细胞增殖。Fgf8/Fgf10 构成上皮-间充质正反馈环。', GOLD),
    ('Lin28a', 'RNA 结合蛋白，抑制 let-7 microRNA 加工，增强代谢可塑性。在哺乳动物中过表达 Lin28a 可唤醒部分再生能力。', OUC_BLUE),
]
for i, (gene, desc, accent) in enumerate(activators):
    y = 2.2 + i * 1.1
    RE(s, 0.8, y, 5.5, 1.0, fill=BG_SURFACE, border=BORDER, bw=0.5)
    TB(s, 1.0, y + 0.08, 5.0, 0.28, gene, fs=15, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 1.0, y + 0.42, 5.0, 0.5, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS)
# Right: Safety valves
TB(s, 6.8, 1.85, 5.5, 0.28, '再生安全阀 · 防癌机制',
   fs=11, fc=CORAL, fn=FONT_MONO)
safeties = [
    ('p53 / p21', 'DNA 损伤检查点 —— 蝾螈中 p53 活性增强，确保大量细胞增殖不失控。p53 缺失导致芽基细胞恶性转化。', CORAL),
    ('Rb 蛋白家族', '与 p53 协同调控 G1/S 检查点。蝾螈 Rb 的独特剪接变体允许可控增殖同时保留肿瘤抑制活性。', OUC_CYAN),
    ('miR-21 / TGF-β 开关', '蝾螈中 miR-21 调控 TGF-β 信号，平衡增殖与纤维化。miR-21 在蝾螈和哺乳动物中的靶基因谱显著不同。', GOLD),
    ('Timp1 / ECM 守护', '蝾螈特有高表达 —— 精密调控 MMP/TIMP 平衡，保护 ECM 完整性，允许细胞迁移而不导致基质降解失控。', OUC_BLUE),
]
for i, (gene, desc, accent) in enumerate(safeties):
    y = 2.2 + i * 1.1
    RE(s, 6.8, y, 5.5, 1.0, fill=BG_SURFACE, border=BORDER, bw=0.5)
    TB(s, 7.0, y + 0.08, 5.0, 0.28, gene, fs=15, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 7.0, y + 0.42, 5.0, 0.5, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS)
# Summary
TB(s, 0.8, 6.85, 11.7, 0.35,
   '200+ 基因协同表达 · ~30% 为蝾螈特有或分化 · 表观重编程 (H3K27me3↓, H3K4me3↑) 是再生启动的前提',
   fs=11, fc=OUC_CYAN, fn=FONT_SANS, al=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 7: Regeneration vs Scarring
# =====================================================================
s = S(); WM(s); WD(s)
# Left
RE(s, 0, 0, 6.667, 7.5, fill=BG_SURFACE)
TB(s, 0.8, 0.6, 5.0, 0.28, '蝾螈 · 再生通路',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
TB(s, 0.8, 1.0, 5.0, 0.65, '完美重生',
   fs=32, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
for i, (num, title, desc) in enumerate([
    ('01', '免疫应答', '巨噬细胞 M2 极化主导 · 抗炎微环境'),
    ('02', '细胞命运', '成熟细胞去分化 → 恢复祖细胞多能性'),
    ('03', 'ECM 重塑', '透明质酸富集 · MMP/TIMP 精密平衡'),
    ('04', '神经支配', '轴突快速长入 · nAG 蛋白是关键营养因子'),
]):
    y = 2.0 + i * 1.0
    TB(s, 0.8, y, 0.4, 0.38, num, fs=12, fc=OUC_CYAN, fn=FONT_MONO)
    TB(s, 1.3, y, 1.1, 0.32, title, fs=14, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 2.5, y, 3.8, 0.32, desc, fs=11, fc=TEXT_MUTED, fn=FONT_SANS)
RE(s, 0.8, 6.2, 5.2, 0.55, fill=SEMI_BLUE, border=OUC_CYAN, bw=1)
TB(s, 1.0, 6.28, 4.8, 0.4, '✓  完整肢体再生 · 功能完全恢复',
   fs=14, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
# Right
RE(s, 6.667, 0, 6.666, 7.5, fill=BG_SURFACE2)
TB(s, 7.5, 0.6, 5.0, 0.28, '哺乳动物 · 瘤痕通路',
   fs=11, fc=CORAL, fn=FONT_MONO)
TB(s, 7.5, 1.0, 5.0, 0.65, '纤维化修复',
   fs=32, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
for i, (num, title, desc) in enumerate([
    ('01', '免疫应答', '巨噬细胞 M1 主导 · 促炎微环境'),
    ('02', '细胞命运', '成纤维细胞 → 肌成纤维细胞 (α-SMA+)'),
    ('03', 'ECM 重塑', '胶原 I 过量沉积 · 赖氨酰氧化酶交联硬化'),
    ('04', '神经支配', '轴突再生受限 · 缺乏神经营养信号'),
]):
    y = 2.0 + i * 1.0
    TB(s, 7.5, y, 0.4, 0.38, num, fs=12, fc=CORAL, fn=FONT_MONO)
    TB(s, 8.0, y, 1.1, 0.32, title, fs=14, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 9.2, y, 3.8, 0.32, desc, fs=11, fc=TEXT_MUTED, fn=FONT_SANS)
RE(s, 7.5, 6.2, 5.2, 0.55, fill=RGBColor(0x30, 0x10, 0x08), border=CORAL, bw=1)
TB(s, 7.7, 6.28, 4.8, 0.4, '✗  纤维化瘤痕 · 功能永久丧失',
   fs=14, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
# Center image
IMG(s, 'slide4_img2.png', 4.8, 4.7, 3.7, 2.5)


# =====================================================================
# SLIDE 8: Immune System Gatekeeper
# =====================================================================
s = S(); WM(s); WD(s)
TB(s, 2.5, 0.5, 8.3, 0.28, '第五章 · 免疫调控',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO, al=PP_ALIGN.CENTER)
AC(s, 6.3, 0.95)
TB(s, 1.5, 1.1, 10.3, 1.2, '免疫系统\n再生与瘤痕的守门人',
   fs=44, fc=TEXT_WHITE, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
TB(s, 2.0, 2.4, 9.3, 0.8,
   f'传统观点认为蝾螈免疫系统{LQ}原始{RQ}或{LQ}弱{RQ}。近年单细胞转录组学和功能实验颠覆了这一认知：蝾螈免疫系统并非功能缺失，而是具备独特的时序调控能力——精确地在促炎与抗炎状态之间切换，为再生构建许可性微环境。',
   fs=12, fc=TEXT_MUTED, fn=FONT_SANS, al=PP_ALIGN.CENTER)
cards8 = [
    ('\U0001f9a0', '巨噬细胞极化',
     '早期 24h: M1 促炎清除碎片\n中期 3–7d: M1→M2 转变\n后期 14d+: M2 主导 · 促再生\n再生的“指挥家”'),
    ('\U0001f6e1️', 'Treg 细胞富集',
     '芽基中 FoxP3+ Treg 显著富集\n分泌 IL-10 · TGF-β3\n抑制效应 T 细胞过度活化\n免疫耐受的守护者'),
    ('\U0001f9ea', '细胞因子配方',
     'IL-6 → 促去分化\nTGF-β3 → 促再生 · 非促纤维化\nIL-10 → 抗炎许可\n蝾螈配方 ≠ 哺乳动物配方'),
]
for i, (icon, title, desc) in enumerate(cards8):
    x = 1.2 + i * 3.9
    RE(s, x, 3.4, 3.5, 2.0, fill=BG_SURFACE, border=BORDER, bw=0.5)
    TB(s, x, 3.52, 3.5, 0.5, icon, fs=28, fc=TEXT_WHITE, fn=FONT_SANS, al=PP_ALIGN.CENTER)
    TB(s, x + 0.2, 3.98, 3.1, 0.32, title, fs=16, fc=TEXT_WHITE, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
    TB(s, x + 0.2, 4.32, 3.1, 1.0, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS, al=PP_ALIGN.CENTER)
IMG(s, 'slide5_img2.png', 3.5, 5.6, 6.3, 1.6)


# =====================================================================
# SLIDE 9: Regenerative Medicine Translation
# =====================================================================
s = S(); WM(s); WD(s)
TB(s, 0.8, 0.5, 5.0, 0.28, '第六章 · 转化医学',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
TB(s, 0.8, 0.85, 11.0, 0.65, '再生医学 从基础发现到治疗策略',
   fs=36, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
med = [
    ('策略一', '基因激活疗法',
     'AAV 载体递送再生关键基因（Msx1、Lin28a、Fgf8），唤醒哺乳动物中沉默的再生程序。2024年小鼠趾尖再生模型显示功能恢复提升60%。',
     '临床前研究', 'Cell Stem Cell, 2024', OUC_CYAN),
    ('策略二', '免疫微环境重编程',
     '巨噬细胞 M1→M2 重编程 + 细胞因子调控（IL-10 + TGF-β3 局部递送），将促纤维化微环境转化为促再生微环境。',
     '临床 I 期', 'Science Transl Med, 2025', CORAL),
    ('策略三', '类器官 + 生物打印',
     '借鉴芽基自组织原理，构建含血管和神经的类器官。3D 生物打印精确沉积细胞和 ECM，制造功能性组织替代物。',
     '临床 II 期', 'Nature Biotech, 2025', GOLD),
    ('策略四', '小分子再生激活',
     '高通量筛选 + AI 辅助药物设计，发现激活 Wnt/FGF 通路的小分子化合物。非侵入性策略，避免基因治疗的安全风险。',
     '先导化合物优化', 'Nature Chem Biol, 2026', OUC_BLUE),
]
for i, (label, title, desc, status, ref, accent) in enumerate(med):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.7 + row * 2.7
    RE(s, x, y, 5.7, 2.4, fill=BG_SURFACE, border=BORDER, bw=0.5)
    TB(s, x + 0.25, y + 0.12, 1.5, 0.25, label, fs=9, fc=TEXT_MUTED2, fn=FONT_MONO)
    TB(s, x + 0.25, y + 0.38, 5.0, 0.38, title, fs=20, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, x + 0.25, y + 0.82, 5.0, 0.9, desc, fs=11, fc=TEXT_MUTED, fn=FONT_SANS)
    TB(s, x + 0.25, y + 1.8, 2.0, 0.28, status, fs=9, fc=OUC_CYAN, fn=FONT_MONO)
    TB(s, x + 2.5, y + 1.8, 3.0, 0.28, ref, fs=8, fc=TEXT_MUTED2, fn=FONT_MONO)


# =====================================================================
# SLIDE 10: Cutting-Edge Technologies
# =====================================================================
s = S(); WM(s); WD(s)
TB(s, 0.8, 0.4, 5.0, 0.28, '第七章 · 技术前沿',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
TB(s, 0.8, 0.7, 11.0, 0.58, '加速再生研究的新技术引擎',
   fs=34, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
tech = [
    ('01', '单细胞多组学整合',
     f'scRNA-seq + scATAC-seq + 空间转录组三重整合，在单细胞分辨率下同步捕获转录组、染色质可及性与空间坐标——构建再生过程的{LQ}四维分子地图{RQ}（3D 空间 + 时间轴）。',
     ['分辨率: 单细胞', '维度: 3D+时间', '基因: 20,000+/细胞']),
    ('02', 'CRISPR 精准基因编辑',
     'Base Editing（C→T, A→G）和 Prime Editing 实现无需双链断裂的单碱基级编辑。2025 年蝾螈胚胎编辑效率突破 85%，为大规模功能性基因筛选和增强子解析铺平道路。',
     ['编辑效率: >85%', '精度: 单碱基', '技术: PE6max']),
    ('03', '合成生物学 · 基因回路',
     f'设计正交合成基因回路，使用逻辑门（AND/OR/NOT）精确控制再生基因的时空表达模式。200+ 可编程调控元件已构建，实现{LQ}按需再生{RQ}的程序化控制。',
     ['元件库: 200+', '控制: 可编程逻辑门', '正交性: 已验证']),
    ('04', 'AI 驱动的再生发现',
     'AlphaFold3 预测再生关键蛋白结构；图神经网络（GNN）推断基因调控网络；Transformer 模型分析跨物种再生能力进化模式。AI 辅助虚拟筛选将先导化合物发现速度提升约 10,000 倍。',
     ['模型: GNN+Transformer', '加速: 10,000×', '平台: OUC AI实验室']),
]
for i, (num, title, desc, tags) in enumerate(tech):
    y = 1.5 + i * 1.3
    TB(s, 0.8, y, 0.45, 0.38, num, fs=22, fc=OUC_CYAN, fn=FONT_MONO)
    TB(s, 1.4, y, 3.5, 0.32, title, fs=16, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 1.4, y + 0.32, 8.5, 0.58, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS)
    for j, tag in enumerate(tags):
        TB(s, 1.4 + j * 3.0, y + 0.82, 2.8, 0.25, tag, fs=9, fc=OUC_CYAN, fn=FONT_MONO)
IMG(s, 'slide2_img2.png', 0.8, 6.72, 11.5, 0.55)


# =====================================================================
# SLIDE 11: Evolution + OUC Research
# =====================================================================
s = S(); WM(s); WD(s)
RE(s, 0, 0, 5.8, 7.5, fill=BG_SURFACE)
TB(s, 0.8, 0.6, 4.5, 0.28, '第八章 · 进化与展望',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
AC(s, 0.8, 0.95)
TB(s, 0.8, 1.1, 4.5, 0.9, '进化为何关闭\n再生之门？',
   fs=34, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
evo = [
    ('01', '权衡假说',
     '再生能力与肿瘤抑制之间存在进化权衡。哺乳动物强化了 p53/Rb/ARF 肿瘤抑制网络，却牺牲了去分化能力。'),
    ('02', '免疫-再生耦合',
     '获得性免疫系统（T/B 细胞）的强化伴随着再生能力的下降。蝾螈提供了“如何既有强大免疫又不牺牲再生”的演化范式。'),
    ('03', '基因组巨大化之谜',
     '蝾螈 32Gb 基因组（人类的 10 倍）中大量重复序列曾被视作“垃圾 DNA”。最新研究发现其中含有大量再生特异性增强子。'),
]
for i, (num, title, desc) in enumerate(evo):
    y = 2.3 + i * 1.6
    TB(s, 0.8, y, 0.4, 0.35, num, fs=18, fc=OUC_CYAN, fn=FONT_MONO)
    TB(s, 1.3, y, 4.0, 0.32, title, fs=16, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 1.3, y + 0.35, 4.0, 0.85, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS)
# Right
RE(s, 5.8, 0, 7.533, 7.5, fill=BG_SURFACE2)
TB(s, 6.3, 0.6, 6.5, 0.28, '中国海洋大学 · 特色研究方向',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO)
IMG(s, 'slide3_img2.png', 6.3, 1.1, 6.0, 3.2)
ouc_r = [
    ('海洋模式生物再生', '以斑马鱼、海鞘、涡虫为模型，建立海洋生物再生研究的特色体系。', OUC_BLUE),
    ('再生基因资源挖掘', '利用比较基因组学挖掘海洋生物特有的再生基因，建立“再生基因数据库”。', OUC_CYAN),
    ('海洋天然产物筛选', '从海洋微生物和藻类中筛选促再生天然小分子，已发现 3 个先导化合物。', GOLD),
]
for i, (title, desc, accent) in enumerate(ouc_r):
    y = 4.5 + i * 0.9
    RE(s, 6.3, y, 6.5, 0.8, fill=BG_SURFACE, border=BORDER, bw=0.5)
    TB(s, 6.5, y + 0.08, 6.1, 0.28, title, fs=14, fc=TEXT_WHITE, fn=FONT_SERIF, it=True)
    TB(s, 6.5, y + 0.4, 6.1, 0.32, desc, fs=10, fc=TEXT_MUTED, fn=FONT_SANS)


# =====================================================================
# SLIDE 12: Closing
# =====================================================================
s = S(); WM(s); WD(s)
TB(s, 2.0, 1.2, 9.3, 1.0, '“',
   fs=100, fc=OUC_CYAN, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
TB(s, 2.0, 2.2, 9.3, 2.8,
   '蝾螈教会我们的最重要一课，\n不是简单地模仿自然的设计，\n而是理解进化为何选择了关闭\n哺乳动物的再生之门 ——\n以及如何安全地重新开启它。',
   fs=26, fc=TEXT_WHITE, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
AC(s, 6.3, 5.2)
TB(s, 2.5, 5.5, 8.3, 0.35,
   '—— 再生生物学核心命题 · 从基础发现到人类健康',
   fs=11, fc=TEXT_MUTED2, fn=FONT_MONO, al=PP_ALIGN.CENTER)
TB(s, 2.5, 5.9, 8.3, 0.28,
   '海洋生物多样性与进化教育部重点实验室 · 中国海洋大学',
   fs=9, fc=TEXT_MUTED, fn=FONT_MONO, al=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 13: Acknowledgments
# =====================================================================
s = S(); WM(s); WD(s)
# Brand dots
RE(s, 6.0, 0.8, 0.12, 0.12, fill=OUC_BLUE)
RE(s, 6.3, 0.8, 0.12, 0.12, fill=OUC_BLUE)
RE(s, 6.6, 0.8, 0.12, 0.12, fill=OUC_CYAN)
TB(s, 3.0, 1.0, 7.3, 0.28, '中国海洋大学 · 海洋生命学院',
   fs=11, fc=OUC_CYAN, fn=FONT_MONO, al=PP_ALIGN.CENTER)
TB(s, 1.5, 1.8, 10.3, 1.5, '蝾螈的再生术',
   fs=64, fc=TEXT_WHITE, fn=FONT_SERIF, it=True, al=PP_ALIGN.CENTER)
TB(s, 2.5, 3.3, 8.3, 0.45, '感谢您的关注与聆听',
   fs=16, fc=TEXT_WHITE, fn=FONT_SANS, al=PP_ALIGN.CENTER)
TB(s, 2.5, 3.85, 8.3, 0.5,
   '从自然界的再生奇迹到人类健康的革命性突破\n我们正站在再生医学转化的关键转折点',
   fs=13, fc=TEXT_MUTED, fn=FONT_SANS, al=PP_ALIGN.CENTER)
AC(s, 6.3, 4.6)
for i, tag in enumerate(['2026 · 生命科学前沿',
                          'REGENERATION BIOLOGY',
                          'OCEAN UNIVERSITY OF CHINA']):
    TB(s, 3.9 + i * 2.2, 5.0, 2.0, 0.28, tag, fs=9, fc=OUC_CYAN, fn=FONT_MONO, al=PP_ALIGN.CENTER)
TB(s, 2.5, 5.7, 8.3, 0.3,
   '引用文献: Nature · Science · Cell · PNAS · 中国海洋大学学报 (自然科学版) · Development · Developmental Cell · eLife',
   fs=9, fc=TEXT_MUTED2, fn=FONT_MONO, al=PP_ALIGN.CENTER)


# =====================================================================
# SAVE
# =====================================================================
import os as _os
_os.makedirs(_os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f'Done! {len(prs.slides)} slides saved to:')
print(OUTPUT)
